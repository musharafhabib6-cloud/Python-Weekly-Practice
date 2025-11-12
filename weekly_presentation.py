from pptx import Presentation

# Create a new PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "Weekly Python Practice Summary"
slide1.placeholders[1].text = "Debugging, Python Basics, NumPy & Pandas, OOP, and Mini Project"

# Slide 2: Debugging & Configuration
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "Debugging & Configuration"
slide2.placeholders[1].text = (
    "• Wrote a small program with a logical bug.\n"
    "• Used debugger to inspect variables and step through code.\n"
    "• Learned to identify and fix logical issues efficiently."
)

# Slide 3: Python Basics
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "Python Basics"
slide3.placeholders[1].text = (
    "• Practiced variables, loops, functions, and lists.\n"
    "• Created simple calculator and grade tracker.\n"
    "• Strengthened understanding of Python syntax and control flow."
)

# Slide 4: NumPy & Pandas
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
slide4.shapes.title.text = "NumPy & Pandas"
slide4.placeholders[1].text = (
    "• Loaded 'grades.csv' dataset using Pandas.\n"
    "• Calculated averages, sorted, and filtered data.\n"
    "• Learned structured data handling for real-world applications."
)

# Slide 5: Practical OOP
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
slide5.shapes.title.text = "Practical OOP (Functions & Classes)"
slide5.placeholders[1].text = (
    "• Refactored code into a Student class.\n"
    "• Implemented methods like calculate_average().\n"
    "• Gained experience in modular and reusable programming."
)

# Slide 6: Mini Project
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
slide6.shapes.title.text = "Mini Project Checkpoint"
slide6.placeholders[1].text = (
    "• Combined all skills into a grade analysis project.\n"
    "• Processed CSV data and generated output.csv.\n"
    "• Completed first standalone Python project successfully."
)

# Slide 7: Summary
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
slide7.shapes.title.text = "Summary"
slide7.placeholders[1].text = (
    "• Practiced debugging, Python basics, data handling, and OOP.\n"
    "• Created a working mini project integrating all concepts.\n"
    "• Built a solid foundation for advanced Python projects."
)

# Save the presentation
prs.save("Weekly_Python_Practice_Presentation.pptx")

print("✅ Presentation created successfully: Weekly_Python_Practice_Presentation.pptx")
